"""Validate present.pptx: no junk, reasonable text length, 12 slides."""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "present.pptx"
SLIDE_W = 9144000  # default 16:9 width in EMU
SLIDE_H = 5143500

JUNK = re.compile(
    r"коворкинг|сводн|большие данные|duckdb|gzip|caffeine|"
    r"20 млн|postgresql|147\.45|dashboard|конструктор|prosto|"
    r"analitics|элитная тима|next slide",
    re.I,
)

ALLOWED_OFF_TOPIC = set()  # empty — anything matching JUNK fails


def shape_overflow(shape, slide_w: int, slide_h: int) -> bool:
    if not hasattr(shape, "text_frame"):
        return False
    text = (shape.text or "").strip()
    if not text or len(text) < 80:
        return False
    # rough: very long text in small box
    if len(text) > 380:
        return True
    try:
        if shape.width and shape.left is not None:
            right = shape.left + shape.width
            if right > slide_w + Emu(50000):
                return True
    except Exception:
        pass
    return False


def main() -> int:
    if not PPTX.exists():
        print(f"Missing {PPTX}", file=sys.stderr)
        return 1
    prs = Presentation(str(PPTX))
    errors: list[str] = []
    if len(prs.slides) != 12:
        errors.append(f"Expected 12 slides, got {len(prs.slides)}")

    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            t = shape.text.strip()
            if not t:
                continue
            if JUNK.search(t):
                errors.append(f"Slide {i}: junk text: {t[:80]!r}")
            if shape_overflow(shape, prs.slide_width, prs.slide_height):
                errors.append(f"Slide {i}: long/overflow? ({len(t)} chars): {t[:60]}…")

    if errors:
        print("VALIDATION FAILED:")
        for e in errors:
            print(" -", e)
        return 1
    print(f"OK: {PPTX} — 12 slides, no junk detected")
    return 0


if __name__ == "__main__":
    sys.exit(main())
