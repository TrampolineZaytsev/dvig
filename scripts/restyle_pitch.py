"""Restyle «ДВИГ. Питч.pptx»: site palette, brighter deck, more icons — text unchanged."""
from __future__ import annotations

import re
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "pitch_source.pptx"
OUT = ROOT / "ДВИГ. Питч — стиль.pptx"
OUT_DL = Path(r"c:\Users\Trampoline\Downloads\ДВИГ. Питч — стиль.pptx")

MAGENTA = RGBColor(0xC6, 0x26, 0x9E)
PURPLE = RGBColor(0x80, 0x64, 0xA2)
GLOW = RGBColor(0xDD, 0xBA, 0xFF)
PINK = RGBColor(0xF5, 0xDF, 0xF0)
DEEP = RGBColor(0x0C, 0x08, 0x18)
BLACK = RGBColor(0x08, 0x04, 0x12)

LOGO = ROOT / "public" / "test_dvig-logo.png"
if not LOGO.exists():
    LOGO = ROOT / "public" / "old_dvig-logo.png"

ICONS = [
    ROOT / "scripts" / "_media_probe" / f
    for f in (
        "image5.png",
        "image7.png",
        "image9.png",
        "image10.png",
        "image17.png",
        "image21.png",
        "image22.png",
        "image23.png",
    )
]
ICONS = [p for p in ICONS if p.exists()]

def icons_for(idx: int, count: int) -> list[Path]:
    if not ICONS:
        return []
    out: list[Path] = []
    for j in range(count):
        out.append(ICONS[(idx + j) % len(ICONS)])
    return out


SLIDE_ICON_COUNT: dict[int, int] = {
    1: 5,
    2: 6,
    3: 6,
    4: 4,
    5: 5,
    6: 5,
    7: 5,
    8: 5,
    9: 6,
    10: 6,
    11: 5,
    12: 6,
    13: 5,
    14: 5,
    15: 6,
    16: 5,
    17: 5,
}

TITLE_FIRST_LINES = {
    "Проблема и контекст",
    "Решение и продукт",
    "Целевая аудитория",
    "Стейкхолдеры",
    "Рынок и возможности",
    "Нормы, власть и ответственность",
    "Бизнес-модель ДВИГ",
    "Этапы развития и трекшн",
    "Команда",
    "Риски",
    "Чего не хватает и как закроем",
    "Запрос и следующий шаг",
    "СПАСИБО ЗА ВНИМАНИЕ!",
}


def iter_shapes(container):
    for sh in container.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(sh)
        else:
            yield sh


def text_of(sh) -> str:
    if not sh.has_text_frame:
        return ""
    return sh.text_frame.text or ""


def is_page_num(t: str) -> bool:
    return bool(re.match(r"^\d{1,2}/\d{1,2}$", t.strip()))


def max_font_pt(sh) -> float:
    best = 0.0
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                best = max(best, run.font.size.pt)
    return best


def is_title_text(t: str, sh, slide_max_pt: float) -> bool:
    t = t.strip()
    if not t or is_page_num(t):
        return False
    first = t.split("\n")[0].strip()
    if first in TITLE_FIRST_LINES or t.startswith("ДВИГ"):
        return True
    pt = max_font_pt(sh)
    if slide_max_pt and pt >= slide_max_pt - 0.5:
        return True
    if pt >= 24:
        return True
    return len(first) < 36 and "\n" not in t[:50]


def style_run(run, *, title: bool) -> None:
    run.font.name = "Segoe UI"
    if title:
        run.font.color.rgb = MAGENTA
        run.font.bold = True
        if run.font.size and run.font.size < Pt(28):
            run.font.size = Pt(min(44, int(run.font.size.pt * 1.06)))
    else:
        run.font.color.rgb = GLOW


def slide_max_font(slide) -> float:
    best = 0.0
    for shp in iter_shapes(slide):
        if shp.has_text_frame and text_of(shp).strip():
            best = max(best, max_font_pt(shp))
    return best


def style_text_shape(sh, slide_max_pt: float) -> None:
    t = text_of(sh).strip()
    if not t:
        return
    if is_page_num(t):
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = PURPLE
                run.font.size = Pt(11)
        return
    title = is_title_text(t, sh, slide_max_pt)
    for pi, para in enumerate(sh.text_frame.paragraphs):
        if not para.runs:
            continue
        for run in para.runs:
            style_run(run, title=title and pi == 0)


def style_shape_fill_line(sh) -> None:
    try:
        if sh.fill.type == MSO_FILL_TYPE.SOLID:
            sh.fill.solid()
            sh.fill.fore_color.rgb = PURPLE
            try:
                sh.fill.transparency = 0.35
            except Exception:
                pass
    except Exception:
        pass
    try:
        if sh.line.fill.type == MSO_FILL_TYPE.SOLID:
            sh.line.color.rgb = MAGENTA
            sh.line.width = Pt(2)
    except Exception:
        pass


def apply_background(slide) -> None:
    fill = slide.background.fill
    try:
        fill.gradient()
        fill.gradient_angle = 135.0
        stops = fill.gradient_stops
        stops[0].color.rgb = MAGENTA
        if len(stops) > 2:
            stops[1].color.rgb = PURPLE
            stops[2].color.rgb = BLACK
        else:
            stops[1].color.rgb = BLACK
    except Exception:
        fill.solid()
        fill.fore_color.rgb = DEEP


def add_glow_orbs(slide, sw: int, slide_h: int) -> None:
    specs = [
        (int(sw * 0.78), int(slide_h * 0.08), int(sw * 0.35), int(sw * 0.35), 0.72),
        (int(-sw * 0.12), int(slide_h * 0.55), int(sw * 0.42), int(sw * 0.42), 0.78),
        (int(sw * 0.55), int(slide_h * 0.72), int(sw * 0.28), int(sw * 0.28), 0.82),
    ]
    for i, (left, top, w, h, trans) in enumerate(specs):
        orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, w, h)
        orb.fill.solid()
        orb.fill.fore_color.rgb = MAGENTA if i % 2 == 0 else PURPLE
        try:
            orb.fill.transparency = trans
        except Exception:
            pass
        orb.line.fill.background()
        orb.name = f"dvig_glow_{i}"


def add_accent_shapes(slide, sw: int, sh: int) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, sw, int(sh * 0.012)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = MAGENTA
    bar.line.fill.background()
    bar.name = "dvig_accent_top"

    side = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, int(sw * 0.008), sh
    )
    side.fill.solid()
    side.fill.fore_color.rgb = PURPLE
    side.line.fill.background()
    side.name = "dvig_accent_side"


def add_logo(slide, sw: int, sh: int) -> None:
    if not LOGO.exists():
        return
    w = int(sw * 0.11)
    h = int(w * 0.35)
    left = sw - w - int(sw * 0.03)
    top = int(sh * 0.04)
    pic = slide.shapes.add_picture(str(LOGO), left, top, width=w, height=h)
    pic.name = "dvig_logo_corner"


def add_icons(slide, idx: int, sw: int, slide_h: int) -> None:
    paths = icons_for(idx, SLIDE_ICON_COUNT.get(idx, 4))
    if not paths:
        return
    n = len(paths)
    margin = int(sw * 0.03)
    size = int(sw * 0.07)
    # bottom row
    y = int(slide_h * 0.855)
    gap = int((sw - 2 * margin - n * size) / max(n - 1, 1))
    for i, path in enumerate(paths):
        left = margin + i * (size + gap)
        pic = slide.shapes.add_picture(str(path), left, y, width=size, height=size)
        pic.name = f"dvig_icon_b_{idx}_{i}"
    # left column (smaller)
    col_n = min(3, len(paths))
    csize = int(sw * 0.055)
    for i in range(col_n):
        top = int(slide_h * 0.18) + i * (csize + int(slide_h * 0.04))
        pic = slide.shapes.add_picture(
            str(paths[i]), int(sw * 0.02), top, width=csize, height=csize
        )
        pic.name = f"dvig_icon_l_{idx}_{i}"


def nudge_layout(shape, sw: int, slide_h: int) -> None:
    t = text_of(shape).strip()
    if not t or is_page_num(t) or not shape.has_text_frame:
        return
    if shape.width < int(sw * 0.15):
        return
    if len(t) > 120 and shape.width < int(sw * 0.55):
        shape.width = min(int(sw * 0.62), int(shape.width * 1.12))


def enlarge_slide4(slide, sw: int, sh: int) -> None:
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if len(pics) >= 2:
        pics[0].width = int(sw * 0.42)
        pics[0].height = int(sh * 0.78)
        pics[0].left = int(sw * 0.04)
        pics[0].top = int(sh * 0.1)
        pics[1].width = int(sw * 0.42)
        pics[1].height = int(sh * 0.78)
        pics[1].left = int(sw * 0.5)
        pics[1].top = int(sh * 0.1)


def restyle_slide(slide, idx: int, sw: int, slide_h: int) -> None:
    apply_background(slide)
    add_glow_orbs(slide, sw, slide_h)
    add_accent_shapes(slide, sw, slide_h)
    max_pt = slide_max_font(slide)
    for shp in iter_shapes(slide):
        if shp.has_text_frame:
            style_text_shape(shp, max_pt)
            nudge_layout(shp, sw, slide_h)
        style_shape_fill_line(shp)
    add_logo(slide, sw, slide_h)
    add_icons(slide, idx, sw, slide_h)
    if idx == 4:
        enlarge_slide4(slide, sw, slide_h)
    if idx == 17:
        for shp in iter_shapes(slide):
            t = text_of(shp).strip()
            if "СПАСИБО" in t:
                for para in shp.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(54)
                        run.font.color.rgb = MAGENTA
                        run.font.bold = True


def main() -> int:
    if not SRC.exists():
        print(f"Missing {SRC}", file=sys.stderr)
        return 1
    print(f"Loading {SRC.name}…")
    prs = Presentation(str(SRC))
    sw, slide_h = prs.slide_width, prs.slide_height
    n = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        print(f" slide {i}/{n}")
        restyle_slide(slide, i, sw, slide_h)
    print(f"Saving {OUT.name}…")
    prs.save(str(OUT))
    shutil.copy2(OUT, OUT_DL)
    print(f"Done:\n  {OUT}\n  {OUT_DL}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
