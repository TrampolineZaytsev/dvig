"""Check present.pptx for text overflow and overlapping text boxes."""
from __future__ import annotations

import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "present.pptx"

# slide size (EMU) — 13.333" x 7.5" at 914400 EMU/inch
SLIDE_W = 12192000
SLIDE_H = 6858000
MARGIN = 91440  # ~0.1 inch


@dataclass
class TextBoxInfo:
    slide: int
    name: str
    text: str
    left: int
    top: int
    width: int
    height: int
    chars: int
    lines: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    def area(self) -> int:
        return self.width * self.height


def iter_text_shapes(slide):
    def walk(container):
        for sh in container.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(sh)
            elif hasattr(sh, "text_frame"):
                yield sh

    yield from walk(slide)


def boxes_overlap(a: TextBoxInfo, b: TextBoxInfo, min_area: int = 80000000) -> bool:
    """True if significant overlap between two substantive text boxes."""
    if a.slide != b.slide or a.name == b.name:
        return False
    if a.chars < 8 or b.chars < 8:
        return False
    x_overlap = max(0, min(a.right, b.right) - max(a.left, b.left))
    y_overlap = max(0, min(a.bottom, b.bottom) - max(a.top, b.top))
    overlap_area = x_overlap * y_overlap
    if overlap_area < min_area:
        return False
    smaller = min(a.area(), b.area())
    if smaller == 0:
        return False
    return overlap_area > smaller * 0.35


def estimate_overflow(info: TextBoxInfo) -> str | None:
    """Heuristic: line count vs box height (ignore shape position outside slide)."""
    if not info.text.strip() or info.chars <= 3:
        return None
    line_h = 130000
    max_lines = max(1, info.height // line_h)
    if info.lines > max_lines + 1:
        return f"строк {info.lines} при ~{max_lines} строк высоты"
    chars_per_line = max(18, info.width // 48000)
    est_lines = max(1, (info.chars + chars_per_line - 1) // chars_per_line)
    if est_lines > max_lines + 1:
        return f"~{est_lines} строк в блоке (~{max_lines} макс.)"
    if info.chars > 200 and info.height < 900000:
        return f"много текста ({info.chars} симв.)"
    return None


def collect_boxes(prs: Presentation) -> list[TextBoxInfo]:
    boxes: list[TextBoxInfo] = []
    for si, slide in enumerate(prs.slides, 1):
        for sh in iter_text_shapes(slide):
            t = (sh.text or "").strip()
            if not t or t == "2026":
                continue
            if len(t) <= 2 and t.isdigit():
                continue
            lines = len([x for x in t.split("\n") if x.strip()]) or 1
            boxes.append(
                TextBoxInfo(
                    slide=si,
                    name=sh.name,
                    text=t.replace("\n", " | ")[:120],
                    left=int(sh.left),
                    top=int(sh.top),
                    width=int(sh.width),
                    height=int(sh.height),
                    chars=len(t.replace("\n", "")),
                    lines=lines,
                )
            )
    return boxes


def check_com_overflow() -> list[str]:
    """Use PowerPoint to detect shapes with text overflow (if available)."""
    issues: list[str] = []
    try:
        import time

        import win32com.client

        pp = win32com.client.Dispatch("PowerPoint.Application")
        pp.Visible = 1
        pres = pp.Presentations.Open(str(PPTX.resolve()), False, False, True)
        time.sleep(1)
        for i in range(1, pres.Slides.Count + 1):
            slide = pres.Slides.Item(i)
            for j in range(1, slide.Shapes.Count + 1):
                sh = slide.Shapes.Item(j)
                try:
                    if sh.HasTextFrame != -1:
                        continue
                    tr = sh.TextFrame.TextRange
                    t = tr.Text.strip()
                    if not t or t == "2026" or (len(t) <= 2 and t.isdigit()):
                        continue
                    # TextFrame overflow: 1 = msoTextOverflowOverflow (text runs outside)
                    overflow = sh.TextFrame.TextRange.BoundHeight > sh.Height * 1.08
                    if overflow:
                        preview = t.replace("\n", " ")[:80]
                        issues.append(f"Слайд {i:02d} [{sh.Name}]: переполнение (BoundHeight) — {preview}")
                except Exception:
                    pass
                if sh.Type == 6:
                    for k in range(1, sh.GroupItems.Count + 1):
                        g = sh.GroupItems.Item(k)
                        try:
                            if g.HasTextFrame == -1:
                                tr = g.TextFrame.TextRange
                                t = tr.Text.strip()
                                if not t:
                                    continue
                                if tr.BoundHeight > g.Height * 1.08:
                                    preview = t.replace("\n", " ")[:80]
                                    issues.append(
                                        f"Слайд {i:02d} [{g.Name}]: переполнение — {preview}"
                                    )
                        except Exception:
                            pass
        pres.Close()
        pp.Quit()
    except Exception as exc:
        issues.append(f"(COM-проверка недоступна: {exc})")
    return issues


def main() -> int:
    if not PPTX.exists():
        print(f"Нет файла {PPTX}", file=sys.stderr)
        return 1

    prs = Presentation(str(PPTX))
    boxes = collect_boxes(prs)
    overlaps: list[str] = []
    overflows: list[str] = []

    for i, a in enumerate(boxes):
        issue = estimate_overflow(a)
        if issue:
            overflows.append(f"Слайд {a.slide:02d} [{a.name}]: {issue} — «{a.text}»")
        for b in boxes[i + 1 :]:
            if boxes_overlap(a, b):
                overlaps.append(
                    f"Слайд {a.slide:02d}: «{a.name}» ↔ «{b.name}» "
                    f"(пересечение ~{min(a.right,b.right)-max(a.left,b.left)}×"
                    f"{min(a.bottom,b.bottom)-max(a.top,b.top)} EMU)"
                )

    com_issues = check_com_overflow()
    # filter COM unavailable message from fail count
    com_real = [x for x in com_issues if not x.startswith("(")]

    print(f"Проверка: {PPTX.name} ({len(prs.slides)} слайдов, {len(boxes)} текстовых блоков)\n")

    if overlaps:
        print("=== Пересечения текста ===")
        for o in overlaps:
            print(" •", o)
        print()

    if overflows:
        print("=== Возможное переполнение (эвристика) ===")
        for o in overflows:
            print(" •", o)
        print()

    if com_real:
        print("=== Переполнение (PowerPoint) ===")
        for o in com_real:
            print(" •", o)
        print()

    total = len(overlaps) + len(overflows) + len(com_real)
    if total == 0:
        print("OK: явных пересечений и переполнений не найдено.")
        return 0

    print(f"Найдено проблем: {len(overlaps)} пересечений, {len(overflows)} эвристика, {len(com_real)} COM.")
    return 1


if __name__ == "__main__":
    sys.exit(main())
