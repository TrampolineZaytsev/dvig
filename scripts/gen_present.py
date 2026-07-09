"""Generate present.pptx (12 slides «Метамышь») — purple/black DVIG brand, no template junk."""
from __future__ import annotations

import io
import subprocess
import sys
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "present.pptx"
PDF = ROOT / "present.pdf"
LOGO = ROOT / "public" / "dvig-logo.png"
SITE = "http://localhost:3000"

# Brand (from globals.css)
BLACK = RGBColor(0x0A, 0x06, 0x12)
PURPLE = RGBColor(0x80, 0x64, 0xA2)
MAGENTA = RGBColor(0xC6, 0x26, 0x9E)
GLOW = RGBColor(0xDD, 0xBA, 0xFF)
WHITE = RGBColor(0xF5, 0xF0, 0xFA)
MUTED = RGBColor(0xA8, 0x98, 0xB8)
CARD = RGBColor(0x1A, 0x12, 0x28)

SLIDES_DATA: list[dict] = [
    {
        "layout": "title",
        "title": "ДВИГ",
        "subtitle": "Живые встречи по интересам",
        "tag": "Зайцев · Журавлёв · Полина · Соня · Настя",
        "foot": "Дело → люди → офлайн",
        "url": SITE,
    },
    {
        "layout": "two",
        "num": "02",
        "title": "Проблема",
        "subtitle": "хаос афиши · из поля",
        "left": "Манеж: три источника времени у стойки.",
        "right": "«Этажи»: пять вкладок — решение голосом.\n«Родина»: сайт ≠ stories в очереди.",
        "foot": "JTBD: «пойти не с кем»",
    },
    {
        "layout": "two",
        "num": "03",
        "title": "Инсайт",
        "subtitle": "ось «общение»",
        "left": "Досуг — совместное действие и доверие к месту.",
        "right": "Решение в чате и голосом.\nАфиша — повод для разговора, не скролл.",
        "foot": "не каталог id событий",
    },
    {
        "layout": "body",
        "num": "04",
        "title": "Решение",
        "subtitle": "ДВИГ + spb-events",
        "body": (
            "KudaGo → нормализация → (опц.) ИИ → чат / export\n"
            "ДВИГ: компания на событие · не dating\n"
            f"Веб-демо: {SITE}/app"
        ),
        "foot": "куратор с ответственностью",
    },
    {
        "layout": "quad",
        "num": "05",
        "title": "Рынок",
        "subtitle": "TAM / SAM / SOM · ≥ 1:30",
        "quads": [
            ("TAM", "досуг в городе"),
            ("SAM", "18–28 · СПб/МСК"),
            ("SOM", "1–2 медиа/чата"),
            ("Платит", "организатор"),
        ],
        "left": "Разрыв: dating (свайпы) vs афиша (куда) — нет «с кем».",
        "right": "Сейчас: пилот CLI для редакций.\nГоризонт: freemium + комиссии.",
    },
    {
        "layout": "features",
        "num": "06",
        "title": "Продукт",
        "subtitle": "spb-events · 8 команд + веб-демо",
        "feats": ["search", "info", "analyze", "export", "share", "trending", "prices", "config"],
        "body": "Next.js · афиша KudaGo · мок групп и safety",
        "url": f"{SITE}/app",
    },
    {
        "layout": "body",
        "num": "07",
        "title": "Доверие",
        "subtitle": "алгоритм ≠ доверие",
        "body": (
            "Факт → KudaGo · смысл → ИИ (маркировка)\n"
            "share → решение редактора\n"
            "Ошибка модели ≠ ложь площадки"
        ),
        "foot": "прозрачность источника",
    },
    {
        "layout": "quad",
        "num": "08",
        "title": "Монетизация",
        "subtitle": "не берём за живое общение",
        "quads": [
            ("Freemium", "пользователям"),
            ("Комиссии", "билеты"),
            ("Размещение", "площадкам"),
            ("Пилот", "без выручки"),
        ],
        "foot": "честная реклама · без скрытого промо",
    },
    {
        "layout": "two",
        "num": "09",
        "title": "Конкуренты",
        "subtitle": "ниша ДВИГ",
        "left": "Dating: Twinby, Bumble, VK Dating — давление «свидания».",
        "right": "Афиши: Timepad, KudaGo — «куда», не «с кем».\nМы: дело + люди + данные.",
        "foot": "KudaGo = источник, не враг",
    },
    {
        "layout": "quad",
        "num": "10",
        "title": "Трекция",
        "subtitle": "сделано и пилот",
        "quads": [
            ("6 записей", "дневника"),
            ("spb-events", "CLI MVP"),
            ("Веб-демо", "/app"),
            ("Метрика", "доверие · мин/пост"),
        ],
        "body": "План: 1–2 студ. медиа СПб · «доверяем / нет»",
    },
    {
        "layout": "meta",
        "num": "11",
        "title": "Метарамка",
        "subtitle": "нормы · власть · ответственность",
        "quads": [
            ("Не видим", "без CLI · смартфона"),
            ("Страдают", "копипаст · непрозрачная реклама"),
            ("Выход", "конфиг · удаление аккаунта"),
            ("Алина", "группа · публичный формат"),
        ],
        "body": "Не «алгоритм защитит». Честный ИИ · слепые зоны · забота.",
    },
    {
        "layout": "title",
        "title": "ДВИГ",
        "subtitle": "Спасибо",
        "tag": "запрос к аудитории",
        "foot": "доверие — через пилот",
        "body": "Пилот: 1–2 студ. медиа или чата в СПб.\nНаставник по safety офлайн.",
        "url": SITE,
    },
]


def set_run_font(run, size: int = 18, bold: bool = False, color: RGBColor = WHITE) -> None:
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size, bold, color)
    return box


def fill_bg(slide, prs) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK


def add_accent_bar(slide, prs) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.07),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = MAGENTA
    bar.line.fill.background()
    glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0.07),
        prs.slide_width,
        Inches(0.03),
    )
    glow.fill.solid()
    glow.fill.fore_color.rgb = PURPLE
    glow.line.fill.background()


def add_glow_orbs(slide, prs) -> None:
    """Soft radial accents like dvig-page on the site."""
    for left, top, size, color in [
        (Inches(8.5), Inches(-0.5), Inches(4.5), PURPLE),
        (Inches(-1.0), Inches(4.5), Inches(3.5), MAGENTA),
    ]:
        orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, size, size)
        orb.fill.solid()
        orb.fill.fore_color.rgb = color
        orb.fill.transparency = 0.82
        orb.line.fill.background()


def add_page_num(slide, num: str) -> None:
    add_textbox(
        slide,
        Inches(12.2),
        Inches(6.85),
        Inches(0.8),
        Inches(0.35),
        num,
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_logo(slide) -> None:
    if not LOGO.exists():
        return
    slide.shapes.add_picture(str(LOGO), Inches(0.45), Inches(0.35), height=Inches(0.55))


def add_card(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    heading: str = "",
    body: str = "",
) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = MAGENTA
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    if heading and body:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = heading
        set_run_font(r, 14, True, GLOW)
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        set_run_font(r2, 12, False, MUTED)
    else:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text or heading
        set_run_font(r, 14, False, WHITE)


def build_title_slide(slide, prs, d: dict) -> None:
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(8), Inches(1.2), d["title"], 54, True, GLOW)
    add_textbox(slide, Inches(0.5), Inches(3.15), Inches(9), Inches(0.7), d["subtitle"], 24, False, GLOW)
    if d.get("tag"):
        add_textbox(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(0.5), d["tag"], 16, False, MUTED)
    if d.get("body"):
        add_textbox(slide, Inches(0.5), Inches(4.7), Inches(9), Inches(1.2), d["body"], 15, False, WHITE)
    if d.get("foot"):
        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.9), Inches(3.2), Inches(0.45)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(0x30, 0x20, 0x48)
        badge.line.color.rgb = PURPLE
        tf = badge.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = d["foot"]
        set_run_font(r, 13, True, GLOW)
    if d.get("url"):
        add_textbox(slide, Inches(0.5), Inches(6.5), Inches(8), Inches(0.4), d["url"], 12, False, PURPLE)


def build_two_slide(slide, prs, d: dict) -> None:
    add_textbox(slide, Inches(0.5), Inches(0.95), Inches(8), Inches(0.8), d["title"], 36, True)
    add_textbox(slide, Inches(0.5), Inches(1.65), Inches(8), Inches(0.45), d["subtitle"], 18, False, GLOW)
    w, h = Inches(5.85), Inches(2.35)
    add_card(slide, Inches(0.5), Inches(2.35), w, h, d["left"])
    add_card(slide, Inches(6.55), Inches(2.35), w, h, d["right"])
    if d.get("foot"):
        add_textbox(slide, Inches(0.5), Inches(5.0), Inches(11), Inches(0.4), d["foot"], 14, True, MAGENTA)


def build_body_slide(slide, prs, d: dict) -> None:
    add_textbox(slide, Inches(0.5), Inches(0.95), Inches(8), Inches(0.8), d["title"], 36, True)
    add_textbox(slide, Inches(0.5), Inches(1.65), Inches(8), Inches(0.45), d["subtitle"], 18, False, GLOW)
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(11.8), Inches(2.8)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = PURPLE
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16)
    tf.margin_top = Pt(14)
    for i, line in enumerate(d["body"].split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = line
        set_run_font(r, 17, False, WHITE)
    if d.get("foot"):
        add_textbox(slide, Inches(0.5), Inches(5.25), Inches(11), Inches(0.4), d["foot"], 14, True, MAGENTA)


def build_quad_slide(slide, prs, d: dict) -> None:
    add_textbox(slide, Inches(0.5), Inches(0.95), Inches(8), Inches(0.8), d["title"], 36, True)
    add_textbox(slide, Inches(0.5), Inches(1.65), Inches(8), Inches(0.45), d["subtitle"], 18, False, GLOW)
    positions = [
        (Inches(0.5), Inches(2.2)),
        (Inches(6.2), Inches(2.2)),
        (Inches(0.5), Inches(3.75)),
        (Inches(6.2), Inches(3.75)),
    ]
    for (left, top), (t, b) in zip(positions, d["quads"]):
        add_card(slide, left, top, Inches(5.5), Inches(1.35), "", heading=t, body=b)
    y = Inches(5.35)
    if d.get("left"):
        add_textbox(slide, Inches(0.5), y, Inches(5.5), Inches(0.9), d["left"], 13, False, MUTED)
    if d.get("right"):
        add_textbox(slide, Inches(6.2), y, Inches(5.5), Inches(0.9), d["right"], 13, False, MUTED)
    if d.get("body") and not d.get("left"):
        add_textbox(slide, Inches(0.5), y, Inches(11.5), Inches(0.6), d["body"], 14, False, WHITE)
    if d.get("foot"):
        add_textbox(slide, Inches(0.5), Inches(6.15), Inches(11), Inches(0.35), d["foot"], 13, True, MAGENTA)


def build_features_slide(slide, prs, d: dict) -> None:
    add_textbox(slide, Inches(0.5), Inches(0.95), Inches(8), Inches(0.8), d["title"], 36, True)
    add_textbox(slide, Inches(0.5), Inches(1.65), Inches(8), Inches(0.45), d["subtitle"], 18, False, GLOW)
    x0, y0 = Inches(0.5), Inches(2.15)
    cw, ch, gap = Inches(2.75), Inches(0.72), Inches(0.18)
    for i, feat in enumerate(d["feats"]):
        col, row = i % 4, i // 4
        left = x0 + col * (cw + gap)
        top = y0 + row * (ch + gap)
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, cw, ch)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x28, 0x1C, 0x40)
        box.line.color.rgb = MAGENTA
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = feat
        set_run_font(r, 13, True, GLOW)
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(11), Inches(0.5), d["body"], 15, False, WHITE)
    if d.get("url"):
        add_textbox(slide, Inches(0.5), Inches(5.45), Inches(11), Inches(0.4), d["url"], 13, False, PURPLE)


def build_meta_slide(slide, prs, d: dict) -> None:
    build_quad_slide(slide, prs, d)
    # highlight frame
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(2.05), Inches(12.0), Inches(3.25)
    )
    frame.fill.background()
    frame.line.color.rgb = MAGENTA
    frame.line.width = Pt(2.5)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    builders = {
        "title": build_title_slide,
        "two": build_two_slide,
        "body": build_body_slide,
        "quad": build_quad_slide,
        "features": build_features_slide,
        "meta": build_meta_slide,
    }

    for i, data in enumerate(SLIDES_DATA):
        slide = prs.slides.add_slide(blank)
        fill_bg(slide, prs)
        add_glow_orbs(slide, prs)
        add_accent_bar(slide, prs)
        add_logo(slide)
        builders[data["layout"]](slide, prs, data)
        num = data.get("num", f"{i + 1:02d}")
        add_page_num(slide, num)

    return prs


def export_pdf() -> None:
    try:
        import time

        import win32com.client

        subprocess.run(
            ["powershell", "-Command", "Get-Process POWERPNT -EA SilentlyContinue | Stop-Process -Force"],
            check=False,
        )
        time.sleep(2)
        pp = win32com.client.Dispatch("PowerPoint.Application")
        pp.Visible = 1
        pres = pp.Presentations.Open(str(OUT.resolve()))
        pres.SaveAs(str(PDF.resolve()), 32)
        pres.Close()
        pp.Quit()
        print(f"PDF: {PDF}")
    except Exception as exc:
        print(f"PDF skip: {exc}", file=sys.stderr)


def main() -> int:
    prs = build_presentation()
    prs.save(str(OUT))
    print(f"PPTX: {OUT} ({len(prs.slides)} slides)")
    export_pdf()
    # validate
    r = subprocess.run([sys.executable, str(ROOT / "scripts" / "validate_present.py")], cwd=str(ROOT))
    return r.returncode


if __name__ == "__main__":
    sys.exit(main())
