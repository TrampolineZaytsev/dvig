"""Build present.pptx from Purple Black template — gradient backgrounds preserved."""
from __future__ import annotations

import io
import re
import shutil
import subprocess
import sys
import time
import zipfile
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Purple Black Modern Marketing Plan Presentation (2).pptx"
OUT = ROOT / "present.pptx"
PDF = ROOT / "present.pdf"
LOGO = ROOT / "public" / "dvig-logo.png"
SITE = "http://localhost:3000"

DUP_ORDER_1BASED = [1, 2, 2, 5, 6, 3, 5, 6, 2, 7, 6, 1]
FILLERS = ("title", "two", "two", "arch", "quad", "feat", "arch", "quad", "two", "traction", "quad", "title")

JUNK_EXACT = {
    "Next Slide", "Коворкинг просто", "Элитная тима", "Большие данные",
    "Сводные таблицы", "Почему мы?", "GZIP", "Caffeine", "cache", "DuckDB",
    "20 млн строк", "Наша команда", "Загрузка", "CSV", "Конструктор",
    "Pivot", "Генерация", "pivot", "AI", "таблиц", "Prosto", "Analitics",
    "Analytics", "Архитектура", "Оптимизация", "Мощные технологии",
    "для сложных задач", "Функциональность", "решения", "Наши", "результаты",
    "Время ответа с кэшированием", "Время ответа без кэширования",
    "Время ответа от AI", "Скорость генерации КЭША",
    "1 сек.", "3,5 мин.", "9 сек.", "4 сек.",
}
JUNK_CONTAINS = (
    "Быстрое создание сводных", "Решение работает ДОЛИ", "20 МИЛЛИОНОВ",
    "147.45.233.15", "dashboards", "КОНСТРУКТОР", "PostgreSQL", "Монолитный",
    "Caffeine", "DuckDB",
)

SLIDES: list[dict] = [
    {
        "title": "ДВИГ",
        "sub": "Найди компанию на событие",
        "tag": "Зайцев · Журавлёв · Полина · Соня · Настя",
        "foot": "Дело → люди → офлайн",
        "url": SITE,
    },
    {
        "title": "Проблема",
        "sub": "афиша есть — компании нет",
        "tag": "наблюдения · СПб",
        "foot": "Манеж · «Этажи» · «Родина»",
        "left": (
            "Пятница: в stories все в компаниях, а в кино или на лекцию "
            "одному некомфортно."
        ),
        "right": (
            "Манеж — три источника времени.\n"
            "«Этажи» — пять вкладок, итог голосом.\n"
            "«Родина» — сайт ≠ stories."
        ),
    },
    {
        "title": "Инсайт",
        "sub": "ось «общение»",
        "tag": "из метаисследования",
        "foot": "не каталог id событий",
        "left": "Досуг — совместное действие и доверие к людям и месту.",
        "right": (
            "«Куда идём» решают в чате и голосом.\n"
            "Афиша — повод для разговора, не скролл."
        ),
    },
    {
        "title": "Решение",
        "sub": "ДВИГ + spb-events",
        "tag": "взлом нормы",
        "foot": "куратор с ответственностью",
        "body": (
            "ДВИГ: интерес → группа 5–15 → офлайн. Не dating — совместное дело.\n\n"
            "spb-events: KudaGo · LLM · Telegram · export — афиша для редакций.\n\n"
            f"Демо: {SITE}/app"
        ),
    },
    {
        "title": "Рынок",
        "sub": "TAM / SAM / SOM",
        "tag": "слайд ≥ 1:30",
        "foot": "B2B2C-пилот · без выручки",
        "q": [
            ("TAM", "городской досуг"),
            ("SAM", "18–28 · СПб/МСК"),
            ("SOM", "1–2 медиа/чата"),
            ("Платит", "организатор"),
        ],
        "left": "Dating — свайпы и «собеседование». Афиши — «куда», не «с кем».",
        "right": "Ниша: повод + компания. Сейчас — пилот. Горизонт — freemium + комиссии.",
    },
    {
        "title": "Продукт",
        "sub": "MVP · 8 команд CLI",
        "tag": "веб-демо Next.js",
        "url": f"{SITE}/app",
        "feats": [
            "search", "info", "analyze", "export",
            "share", "trending", "prices", "config",
        ],
        "body": "Афиша KudaGo · мок групп и safety",
    },
    {
        "title": "Доверие",
        "sub": "алгоритм ≠ доверие",
        "tag": "факт / ИИ / человек",
        "foot": "прозрачность источника",
        "body": (
            "Факт — KudaGo и организатор. Смысл — LLM с маркировкой.\n\n"
            "share — решение редактора. Ошибка модели исправима; "
            "сломанное доверие — если выдали ИИ за официальное."
        ),
    },
    {
        "title": "Монетизация",
        "sub": "не берём за общение",
        "tag": "платит посредник",
        "foot": "честная реклама",
        "q": [
            ("Freemium", "старт бесплатно"),
            ("Комиссии", "билеты"),
            ("Размещение", "площадкам"),
            ("Пилот", "без выручки"),
        ],
    },
    {
        "title": "Конкуренты",
        "sub": "ниша между dating и афишей",
        "tag": "KudaGo = данные",
        "foot": "не дублируем агрегатор",
        "left": "Dating — стресс и витрина анкет. VK/Telegram — не новые связи в реале.",
        "right": "Timepad, Meetup — без социального слоя. Мы: дело + люди + данные.",
    },
    {
        "title": "Трекция",
        "sub": "сделано и пилот",
        "tag": "метаисследование",
        "foot": "доверие важнее CTR",
        "q": [
            ("6 записей", "дневника"),
            ("spb-events", "CLI MVP"),
            ("Веб-демо", "/app"),
            ("Пилот", "1–2 чата"),
        ],
        "body": "Метрика: минуты на пост и «доверяем / нет»",
    },
    {
        "title": "Метарамка",
        "sub": "нормы · власть",
        "tag": "отдельный слайд",
        "foot": "безопасность = забота",
        "q": [
            ("Не видим", "без CLI"),
            ("Страдают", "копипаст-SMM"),
            ("Выход", "удаление аккаунта"),
            ("Алина", "группа · safety"),
        ],
        "body": "Не «алгоритм защитит». Честный ИИ и публичный формат.",
    },
    {
        "title": "ДВИГ",
        "sub": "Спасибо",
        "tag": "запрос к аудитории",
        "foot": "доверие — через пилот",
        "url": SITE,
        "body": "Пилот: 1–2 студ. медиа/чата СПб.\nНаставник по safety или 152-ФЗ.",
    },
]


def duplicate_slide(pres: Presentation, index: int):
    source = pres.slides[index]
    dest = pres.slides.add_slide(pres.slide_layouts[6])
    for sh in source.shapes:
        el = deepcopy(sh.element)
        dest.shapes._spTree.insert_element_before(el, "p:extLst")
    return dest


def build_structure() -> Presentation:
    prs = Presentation(str(TEMPLATE))
    for idx in DUP_ORDER_1BASED:
        duplicate_slide(prs, idx - 1)
    xml_slides = prs.slides._sldIdLst
    for i in range(8, -1, -1):
        sld_id = xml_slides[i]
        prs.part.drop_rel(sld_id.rId)
        del xml_slides[i]
    return prs


def iter_all_shapes(container):
    for sh in container.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all_shapes(sh)
        else:
            yield sh


def find_shape(slide, name: str):
    for sh in iter_all_shapes(slide):
        if sh.name == name:
            return sh
    return None


def set_shape_text(shape, text: str) -> None:
    if shape is None or not hasattr(shape, "text_frame"):
        return
    shape.text_frame.text = text.replace("\n", "\r")


def clear_junk(slide) -> None:
    for sh in iter_all_shapes(slide):
        if not hasattr(sh, "text_frame"):
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        if t in JUNK_EXACT or any(p in t for p in JUNK_CONTAINS):
            set_shape_text(sh, "")
            continue
        if re.match(r"^Page \d+$", t) or "Ссылка на сайт" in t:
            set_shape_text(sh, "")
        if re.match(r"^https?://147", t):
            set_shape_text(sh, "")


def set_page(slide, num: int) -> None:
    for name in ("TextBox 39", "TextBox 28", "TextBox 8", "TextBox 13", "TextBox 4"):
        sh = find_shape(slide, name)
        if sh and sh.text_frame.text.strip().startswith("Page"):
            set_shape_text(sh, f"{num:02d}")
            return
    for sh in iter_all_shapes(slide):
        if sh.text_frame.text.strip() in {f"{num:02d}", f"{num}"}:
            return
        if re.match(r"^Page \d+$", sh.text_frame.text.strip()):
            set_shape_text(sh, f"{num:02d}")
            return


def fill_title(slide, d: dict) -> None:
    set_shape_text(find_shape(slide, "TextBox 17"), d["title"])
    set_shape_text(find_shape(slide, "TextBox 18"), d["sub"])
    set_shape_text(find_shape(slide, "TextBox 19"), d.get("tag", ""))
    if d.get("body"):
        extra = d["body"]
        if d.get("foot"):
            extra += f"\n\n{d['foot']}"
        set_shape_text(find_shape(slide, "TextBox 21"), extra)
    else:
        set_shape_text(find_shape(slide, "TextBox 21"), d.get("foot", ""))
    set_shape_text(find_shape(slide, "TextBox 23"), d.get("url", ""))


def fill_two(slide, d: dict) -> None:
    set_shape_text(find_shape(slide, "TextBox 29"), d["title"])
    set_shape_text(find_shape(slide, "TextBox 30"), d["sub"])
    set_shape_text(find_shape(slide, "TextBox 35"), d.get("tag", ""))
    set_shape_text(find_shape(slide, "TextBox 32"), d.get("left", ""))
    set_shape_text(find_shape(slide, "TextBox 33"), d.get("right", ""))
    set_shape_text(find_shape(slide, "TextBox 34"), d.get("foot", ""))


def fill_arch(slide, d: dict) -> None:
    title_box = find_shape(slide, "TextBox 17") or find_shape(slide, "TextBox 8")
    body_box = find_shape(slide, "TextBox 16")
    set_shape_text(title_box, d["title"])
    set_shape_text(body_box, d.get("body", ""))
    meta = d.get("tag") or d.get("sub") or d.get("foot", "")
    set_shape_text(find_shape(slide, "TextBox 15"), meta)


def fill_quad(slide, d: dict) -> None:
    set_shape_text(find_shape(slide, "TextBox 14"), d["title"])
    set_shape_text(find_shape(slide, "TextBox 13"), d.get("sub", ""))
    for i, (head, sub) in enumerate(d.get("q", [])):
        set_shape_text(find_shape(slide, f"TextBox {16 + i}"), f"{head}\n{sub}")
    if d.get("left"):
        set_shape_text(find_shape(slide, "TextBox 24"), d["left"])
    if d.get("right"):
        set_shape_text(find_shape(slide, "TextBox 28"), d["right"])
    if d.get("body") and not d.get("left"):
        set_shape_text(find_shape(slide, "TextBox 24"), d["body"])


def fill_feat(slide, d: dict) -> None:
    set_shape_text(find_shape(slide, "TextBox 6"), d["title"])
    set_shape_text(find_shape(slide, "TextBox 7"), d["sub"])
    set_shape_text(find_shape(slide, "TextBox 3"), d.get("tag", ""))
    set_shape_text(find_shape(slide, "TextBox 48"), d.get("url", ""))
    set_shape_text(find_shape(slide, "TextBox 8"), d.get("body", ""))
    feat_shapes = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for g in sh.shapes:
                if g.name.startswith("TextBox") and hasattr(g, "text_frame"):
                    feat_shapes.append(g)
    feat_shapes.sort(key=lambda s: (int(s.top), int(s.left)))
    for sh in feat_shapes:
        set_shape_text(sh, "")
    for sh, label in zip(feat_shapes[: len(d.get("feats", []))], d.get("feats", [])):
        set_shape_text(sh, label)


def fill_traction(slide, d: dict) -> None:
    set_shape_text(find_shape(slide, "TextBox 8"), d["title"])
    set_shape_text(find_shape(slide, "TextBox 9"), d["sub"])
    set_shape_text(find_shape(slide, "TextBox 6"), d.get("tag", ""))
    set_shape_text(find_shape(slide, "TextBox 34"), d.get("foot", ""))
    labels = d.get("q", [])
    for name, (head, sub) in zip(
        ("TextBox 13", "TextBox 21", "TextBox 25", "TextBox 17"),
        labels,
    ):
        set_shape_text(find_shape(slide, name), f"{head}\n{sub}")
    for name in ("TextBox 26", "TextBox 35", "TextBox 36"):
        set_shape_text(find_shape(slide, name), "")
    set_shape_text(find_shape(slide, "TextBox 30"), d.get("body", ""))


FILL_MAP = {
    "title": fill_title,
    "two": fill_two,
    "arch": fill_arch,
    "quad": fill_quad,
    "feat": fill_feat,
    "traction": fill_traction,
}


def fill_presentation(prs: Presentation) -> None:
    for i, (kind, data) in enumerate(zip(FILLERS, SLIDES)):
        slide = prs.slides[i]
        clear_junk(slide)
        FILL_MAP[kind](slide, data)
        set_page(slide, i + 1)
        clear_junk(slide)


def replace_logos() -> None:
    if not LOGO.exists():
        return
    from PIL import Image

    backup = ROOT / "present.before-logo.pptx"
    shutil.copy2(OUT, backup)
    replace_media = ("image3.png", "image4.png", "image11.jpeg")

    def resize_logo(size: tuple[int, int]) -> bytes:
        logo = Image.open(LOGO).convert("RGBA")
        logo.thumbnail(size, Image.Resampling.LANCZOS)
        canvas = Image.new("RGBA", size, (0, 0, 0, 0))
        ox = (size[0] - logo.width) // 2
        oy = (size[1] - logo.height) // 2
        canvas.paste(logo, (ox, oy), logo)
        buf = io.BytesIO()
        canvas.save(buf, format="PNG")
        return buf.getvalue()

    with zipfile.ZipFile(OUT, "r") as zin:
        out_buf = io.BytesIO()
        with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for name in zin.namelist():
                data = zin.read(name)
                base = Path(name).name
                if base in replace_media and name.startswith("ppt/media/"):
                    im = Image.open(io.BytesIO(data))
                    data = resize_logo(im.size)
                zout.writestr(name, data)
        OUT.write_bytes(out_buf.getvalue())


def export_pdf() -> None:
    try:
        import win32com.client

        subprocess.run(
            ["powershell", "-Command", "Get-Process POWERPNT -EA SilentlyContinue | Stop-Process -Force"],
            check=False,
        )
        time.sleep(2)
        pp = win32com.client.Dispatch("PowerPoint.Application")
        pp.Visible = 1
        pres = pp.Presentations.Open(str(OUT.resolve()))
        time.sleep(1)
        pres.SaveAs(str(PDF.resolve()), 32)
        pres.Close()
        pp.Quit()
        print(f"PDF: {PDF}")
    except Exception as exc:
        print(f"PDF skip: {exc}", file=sys.stderr)


def main() -> int:
    """Delegate to PowerPoint COM — python-pptx slide copy breaks media (file won't open)."""
    com = ROOT / "scripts" / "gen_present_com.ps1"
    if not com.exists():
        print(f"Missing {com}", file=sys.stderr)
        return 1
    r = subprocess.run(
        ["powershell", "-ExecutionPolicy", "Bypass", "-File", str(com)],
        cwd=str(ROOT),
    )
    if r.returncode != 0:
        return r.returncode
    val = subprocess.run([sys.executable, str(ROOT / "scripts" / "validate_present.py")], cwd=str(ROOT))
    return val.returncode


if __name__ == "__main__":
    sys.exit(main())
